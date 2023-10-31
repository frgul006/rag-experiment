import unittest
from io import BytesIO

from pptx import Presentation
from reportlab.pdfgen import canvas

from regent_rag.core.extractors import extract_text_from_pdf, extract_text_from_pptx


class TestExtractText(unittest.TestCase):
    def test_extract_text_from_pdf(self) -> None:
        # Create a PDF file with known content using reportlab
        pdf_file = BytesIO()
        c = canvas.Canvas(pdf_file)
        text = "Test content"
        c.drawString(100, 750, text)
        c.save()

        # Call the function with the content of the PDF file
        pdf_file.seek(0)
        result = extract_text_from_pdf(pdf_file.read())

        # Assert that the result is equal to the known content
        self.assertEqual(result, text)

    def test_extract_text_from_pptx(self) -> None:
        # Create a PPTX file with known content
        presentation = Presentation()
        slide_layout = presentation.slide_layouts[0]
        slide = presentation.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = "Test content"

        # Save the PPTX file into bytes
        pptx_file = BytesIO()
        presentation.save(pptx_file)
        pptx_file.seek(0)

        # Call the function with the content of the PPTX file
        result = extract_text_from_pptx(pptx_file.read())

        # Assert that the result is equal to the known content
        self.assertEqual(result, "Test content")


if __name__ == "__main__":
    unittest.main()
