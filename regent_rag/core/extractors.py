from io import BytesIO

from pptx import Presentation
from pypdf import PdfReader


def extract_text_from_pdf(file_content: bytes) -> str:
    """
    Extracts text from a PDF file.

    Parameters:
    file_content (bytes): The content of a PDF file.

    Returns:
    str: The extracted text from the PDF file.
    """
    # Load the PDF file
    pdf_file = PdfReader(BytesIO(file_content))

    # Extract text from each page
    extracted_text = " ".join([page.extract_text() for page in pdf_file.pages])

    return extracted_text.strip()


def extract_text_from_pptx(file_content: bytes) -> str:
    """
    Extracts text from a PPTX file.

    Parameters:
    file_content (bytes): The content of a PPTX file.

    Returns:
    str: The extracted text from the PPTX file.
    """
    # Load the PPTX presentation
    presentation = Presentation(BytesIO(file_content))

    # Extract text from each slide, shape, paragraph, and run
    extracted_text = [
        run.text
        for slide in presentation.slides
        for shape in slide.shapes
        if shape.has_text_frame
        for paragraph in shape.text_frame.paragraphs
        for run in paragraph.runs
    ]

    # Join the extracted text into a single string
    text = " ".join(extracted_text)

    return text.strip()
